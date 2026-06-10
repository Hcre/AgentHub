"""PPT 预览端点 GET /api/fs/pptx-slides 测试（三路径 + 边界）。

用 python-pptx 现造一个真实 .pptx，验证文本/备注抽取、非 pptx 415、不存在 404、损坏 422。
"""

from __future__ import annotations

import base64
import os

os.environ.setdefault("SECRET_KEY", base64.b64encode(b"0" * 32).decode("ascii"))
os.environ.setdefault("DATABASE_URL", "sqlite+aiosqlite:///:memory:")
os.environ.setdefault("LLM_ADAPTER_MODE", "mock")
os.environ.setdefault("ENV", "test")

import importlib.util as _ilu

import pytest
from fastapi import FastAPI
from fastapi.testclient import TestClient

_spec = _ilu.spec_from_file_location("_skills_mod", "app/api/routers/skills.py")
_skills_mod = _ilu.module_from_spec(_spec)
_spec.loader.exec_module(_skills_mod)


@pytest.fixture
def client() -> TestClient:
    app = FastAPI()
    app.include_router(_skills_mod.FS_ROUTER)
    return TestClient(app)


def _make_pptx(path: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # slide 1：标题 + 副标题（标题版式）
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "AgentHub 演示"
    s1.placeholders[1].text = "多 Agent 协作平台"
    # slide 2：标题 + 正文 + 备注
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "架构"
    tf = s2.placeholders[1].text_frame
    tf.text = "5 层洋葱"
    box = s2.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1))
    box.text_frame.text = "CLI/SDK 双轨"
    s2.notes_slide.notes_text_frame.text = "讲者备注：强调依赖倒置"
    prs.save(path)


def test_pptx_slides_happy(client: TestClient, tmp_path) -> None:
    p = str(tmp_path / "demo.pptx")
    _make_pptx(p)
    resp = client.get("/api/fs/pptx-slides", params={"path": p})
    assert resp.status_code == 200, resp.text
    data = resp.json()
    assert data["ok"] is True
    assert data["count"] == 2
    s = data["slides"]
    assert s[0]["title"] == "AgentHub 演示"
    assert "多 Agent 协作平台" in s[0]["texts"]
    assert s[1]["title"] == "架构"
    assert "5 层洋葱" in s[1]["texts"]
    assert "CLI/SDK 双轨" in s[1]["texts"]
    assert "讲者备注" in s[1]["notes"]


def test_pptx_slides_non_pptx_415(client: TestClient, tmp_path) -> None:
    p = str(tmp_path / "note.txt")
    with open(p, "w", encoding="utf-8") as f:
        f.write("not a pptx")
    resp = client.get("/api/fs/pptx-slides", params={"path": p})
    assert resp.status_code == 415, resp.text


def test_pptx_slides_missing_404(client: TestClient, tmp_path) -> None:
    resp = client.get("/api/fs/pptx-slides", params={"path": str(tmp_path / "nope.pptx")})
    assert resp.status_code == 404, resp.text


def test_pptx_slides_corrupt_422(client: TestClient, tmp_path) -> None:
    p = str(tmp_path / "broken.pptx")
    with open(p, "wb") as f:
        f.write(b"this is not a zip/pptx package")
    resp = client.get("/api/fs/pptx-slides", params={"path": p})
    assert resp.status_code == 422, resp.text
